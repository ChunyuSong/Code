from nilearn import plotting
import nibabel as nib
import matplotlib.pyplot as plt
import matplotlib as mpl
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import numpy as np

param = ['FiberFA', 'FiberAxial', 'FiberRadial', 'FiberRatio', 'Non-Restricted', 'FA', 'Axial', 'Radial']
stat = ['tstat1','tstat2','tstat2','tstat1','tstat2','tstat1','tstat1','tstat2']
compare = ['Control > Patient',r"Control and Patient No Difference",'Control and Patient\n\tNo Difference','Control > Patient','Control < Patient','Control > Patient','Control and Patient\n\tNo Difference','Control < Patient']
param_id = ['Fractional Anisotropy','Axial Diffusivity','Radial Diffusivity','Fiber Fraction','Isotropic Non-Restricted Fraction','Fractional Anisotropy','Axial Diffusivity','Radial Diffusivity']
prs = Presentation()
fig = plt.figure(figsize=(20, :sunglasses:,facecolor='k')
for k in range(1,9):
   vmin=0
   fig_id = plt.subplot(3,5,k)
   #plt.subplots_adjust(0.12, 0.03, 0.92, 0.9, 0.1, None)
   img = nib.load('Z:\Peng_Sun\common\\atlas\JHU\MNI152_T1_1mm.nii.gz')
   img_data = img.get_data()
   img_data = img_data[80:100,65:170,50:117]
   img = nib.Nifti1Image(img_data,np.eye(4))
   display = plotting.plot_anat(img,cut_coords=[-9],display_mode='x',annotate=False,cmap='gray',vmin=3000,vmax=10000,axes=fig_id)
   img_stat = nib.load('Z:\Peng_Sun\TBSS\PO1\FA_P01_JHU\stats_lesion_PPMS\\temp\\temp\\tbss_mask_fill_tbss_%s_tfce_corrp_%s.nii.gz'%(param[k-1],stat[k-1]))
   img_stat_data = np.nan_to_num(img_stat.get_data())
   img_stat_data = img_stat_data[80:100,65:170,50:117]
   img_stat = nib.Nifti1Image(img_stat_data,np.eye(4))
   display.add_overlay(img_stat,cmap='autumn',threshold=0)
   # plt.title("%s"%param_id[k-1],fontsize=13,weight='bold',color='white',size='large')
   # plt.text(25,-10,"%s"%compare[k-1],weight='bold',color='white')
   # if k == 1:
   #     plt.text(-25,30,"DBSI",weight='bold',color='white',size='x-large')
   # elif k == 6:
   #     plt.text(-25,30,"DTI",weight='bold',color='white',size='x-large')

fig_id = plt.subplot(3,5,10)
img = nib.load('Z:\Peng_Sun\common\\atlas\JHU\MNI152_T1_1mm.nii.gz')
display = plotting.plot_anat(img,cut_coords=[-1],display_mode='x',annotate=False,cmap='gray',vmin=3000,vmax=10000,axes=fig_id)
img_data = img.get_data()
# plt.plot(range(-46,40),np.repeat(-20,86),color='green')
# plt.plot(np.repeat(0,40),range(40,80),color='green')

norm = mpl.colors.Normalize(vmin=0.05,vmax=1)
sm = plt.cm.ScalarMappable(cmap='autumn', norm=norm)
sm.set_array([])
cax = fig.add_axes([0.76, 0.35, 0.15, 0.03])
cax.set_xticks([])
cb = fig.colorbar(sm, cax=cax, orientation="horizontal",)
cb.set_label('p-value',color='white',weight='bold')
cb.outline.set_visible(False)
cb.set_ticks([np.linspace(0.05,1,2,endpoint=True)])
cb.ax.set_xticklabels(['0.05','1'],color='white',weight='bold')
mng = plt.get_current_fig_manager()
mng.window.showMaximized()
plt.savefig(os.path.join('Z:\Peng_Sun\TBSS\PO1\FA_P01_JHU\stats_lesion_PPMS','Plot.tiff'),bbox_inches='tight',facecolor=fig.get_facecolor())
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
left = Inches(0)
top = Inches(1)
width = Inches(10)
height = Inches(6)
slide.shapes.add_picture(os.path.join('Z:\Peng_Sun\TBSS\PO1\FA_P01_JHU\stats_lesion_PPMS','Plot.png'),left, top, width, height)
prs.save(os.path.join('Z:\Peng_Sun\TBSS\PO1\FA_P01_JHU\stats_lesion_PPMS','Image_Plot.pptx'))
os.remove(os.path.join('Z:\Peng_Sun\TBSS\PO1\FA_P01_JHU\stats_lesion_PPMS','Plot.png'))
plt.close('all')